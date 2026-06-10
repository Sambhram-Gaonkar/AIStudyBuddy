from pptx import Presentation


def extract_pptx_pages(file_path):
    presentation = Presentation(file_path)
    pages = []

    for slide_index, slide in enumerate(presentation.slides, start=1):
        text_parts = []
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text:
                text_parts.append(shape.text.strip())
        text = ' '.join(part for part in text_parts if part)
        if text:
            pages.append({'page_number': slide_index, 'text': text})

    return pages
